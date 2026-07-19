"""Parse a PPTX into a `ParsedDocument` (Module A.1).

PowerPoint carries more structure than a PDF: title placeholders, body text
frames, tables, pictures and speaker notes are all explicit. So mapping is
direct rather than heuristic — the title placeholder becomes a top-level
heading, body text frames become paragraphs or lists, tables become table
blocks, pictures become image refs, and each slide's notes become
`Page.notes` (the raw material for narration scripts later).
"""

from __future__ import annotations

from datetime import datetime, timezone
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schema import Block, BlockKind, ImageRef, Page, ParsedDocument, SourceInfo

_PARSER = "python-pptx"


def _parser_version() -> str:
    try:
        return version("python-pptx")
    except PackageNotFoundError:
        return "unknown"


def _text_frame_blocks(text_frame) -> list[Block]:
    """Turn a text frame into blocks.

    A multi-line text frame is treated as a list: on a slide, a body placeholder
    with several lines is almost always a bullet list, and python-pptx does not
    reliably expose bullet formatting (a flat bullet list and flat body text
    both report indent level 0), so the line count is the most dependable
    signal. A single line stays a paragraph.
    """
    paragraphs = [p.text.strip() for p in text_frame.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    if len(paragraphs) >= 2:
        return [Block(kind=BlockKind.list, items=paragraphs)]
    return [Block(kind=BlockKind.paragraph, text=paragraphs[0])]


def _shape_entries(
    shape, slide_index: int, title_id, image_index: int, warnings: list[str]
) -> tuple[list[tuple[int, Block]], int]:
    """Blocks contributed by one shape, plus the running image counter.

    Returns ``(entries, image_index)`` where each entry is ``(top, Block)`` so
    the caller can restore reading order, and ``image_index`` is the updated
    per-slide picture count.
    """
    top = shape.top or 0

    if title_id is not None and shape.shape_id == title_id:
        text = shape.text_frame.text.strip() if shape.has_text_frame else ""
        if text:
            return [(top, Block(kind=BlockKind.heading, text=" ".join(text.split()), level=1))], image_index
        return [], image_index

    if getattr(shape, "has_table", False):
        rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
        return [(top, Block(kind=BlockKind.table, rows=rows))], image_index

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_index += 1
        width = height = None
        try:
            width, height = shape.image.size
        except Exception:  # noqa: BLE001 - image metadata is best-effort
            warnings.append(f"slide {slide_index}: could not read image dimensions")
        ref = ImageRef(id=f"s{slide_index}-img{image_index}", width=width, height=height)
        return [(top, Block(kind=BlockKind.image, image=ref))], image_index

    if shape.has_text_frame:
        return [(top, block) for block in _text_frame_blocks(shape.text_frame)], image_index

    return [], image_index


def _build_slide(slide, slide_index: int, warnings: list[str]) -> Page:
    """Assemble one slide's blocks in reading order, plus speaker notes."""
    title = slide.shapes.title
    title_id = title.shape_id if title is not None else None

    entries: list[tuple[int, Block]] = []
    image_index = 0
    for shape in slide.shapes:
        shape_entries, image_index = _shape_entries(shape, slide_index, title_id, image_index, warnings)
        entries.extend(shape_entries)
    entries.sort(key=lambda item: item[0])

    notes = None
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip() or None

    return Page(index=slide_index, kind="slide", blocks=[block for _, block in entries], notes=notes)


def parse_pptx(path: str | Path) -> ParsedDocument:
    """Read a PPTX and return its structured representation."""
    path = Path(path)
    warnings: list[str] = []
    prs = Presentation(str(path))
    pages = [_build_slide(slide, index, warnings) for index, slide in enumerate(prs.slides, start=1)]

    core = prs.core_properties
    source = SourceInfo(
        filename=path.name,
        format="pptx",
        page_count=len(prs.slides),
        title=core.title or None,
        author=core.author or None,
        created=core.created,
        modified=core.modified,
    )
    return ParsedDocument(
        source=source,
        parser=_PARSER,
        parser_version=_parser_version(),
        parsed_at=datetime.now(timezone.utc),
        pages=pages,
        warnings=warnings,
    )
