"""Parse a PPTX into a `ParsedDocument` (Module A.1).

PowerPoint carries more structure than a PDF: title placeholders, body text
frames, tables, pictures and speaker notes are all explicit. So mapping is
direct rather than heuristic — the title placeholder becomes a top-level
heading, body text frames become paragraphs or lists, tables become table
blocks, pictures become image refs, and each slide's notes become
`Page.notes` (the raw material for narration scripts later).
"""

from __future__ import annotations

from datetime import datetime, timezone
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schema import Block, BlockKind, ImageRef, Page, ParsedDocument, SourceInfo

_PARSER = "python-pptx"


def _parser_version() -> str:
    try:
        return version("python-pptx")
    except PackageNotFoundError:
        return "unknown"


def _text_frame_blocks(text_frame) -> list[Block]:
    """A text frame becomes a list (several lines) or a single paragraph."""
    paragraphs = [p.text.strip() for p in text_frame.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    if len(paragraphs) >= 2:
        return [Block(kind=BlockKind.list, items=paragraphs)]
    return [Block(kind=BlockKind.paragraph, text=paragraphs[0])]


def parse_pptx(path: str | Path) -> ParsedDocument:
    """Read a PPTX and return its structured representation."""
    path = Path(path)
    warnings: list[str] = []
    prs = Presentation(str(path))
    pages: list[Page] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title
        title_id = title.shape_id if title is not None else None
        positioned: list[tuple[int, Block]] = []
        image_count = 0

        for shape in slide.shapes:
            top = shape.top or 0

            if title_id is not None and shape.shape_id == title_id:
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if text:
                    positioned.append((top, Block(kind=BlockKind.heading, text=" ".join(text.split()), level=1)))
                continue

            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                positioned.append((top, Block(kind=BlockKind.table, rows=rows)))
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                width = height = None
                try:
                    width, height = shape.image.size
                except Exception:  # noqa: BLE001 - image metadata is best-effort
                    warnings.append(f"slide {slide_index}: could not read image dimensions")
                positioned.append((
                    top,
                    Block(kind=BlockKind.image, image=ImageRef(id=f"s{slide_index}-img{image_count}", width=width, height=height)),
                ))
                continue

            if shape.has_text_frame:
                for block in _text_frame_blocks(shape.text_frame):
                    positioned.append((top, block))

        positioned.sort(key=lambda item: item[0])

        notes = None
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            notes = notes_text or None

        pages.append(Page(index=slide_index, kind="slide", blocks=[block for _, block in positioned], notes=notes))

    core = prs.core_properties
    source = SourceInfo(
        filename=path.name,
        format="pptx",
        page_count=len(prs.slides),
        title=core.title or None,
        author=core.author or None,
        created=core.created,
        modified=core.modified,
    )
    return ParsedDocument(
        source=source,
        parser=_PARSER,
        parser_version=_parser_version(),
        parsed_at=datetime.now(timezone.utc),
        pages=pages,
        warnings=warnings,
    )
