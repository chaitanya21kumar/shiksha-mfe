"""Tests for the PPTX parser.

The fixture builds a small but representative deck at runtime (title, a bullet
list, an image and speaker notes on slide 1; a title and a table on slide 2)
so the test is real and self-contained.
"""

import io

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.ingestion.pptx_parser import parse_pptx
from app.ingestion.schema import BlockKind


@pytest.fixture
def sample_pptx(tmp_path):
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    s1.shapes.title.text = "Photosynthesis"
    body = s1.placeholders[1].text_frame
    body.text = "Sunlight"
    body.add_paragraph().text = "Water"
    body.add_paragraph().text = "Carbon dioxide"
    png = io.BytesIO()
    Image.new("RGB", (80, 60), (120, 160, 200)).save(png, "PNG")
    png.seek(0)
    s1.shapes.add_picture(png, Inches(1), Inches(4), Inches(1), Inches(0.75))
    s1.notes_slide.notes_text_frame.text = "Explain the three inputs to photosynthesis."

    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s2.shapes.title.text = "Key terms"
    table = s2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Meaning"
    table.cell(1, 0).text = "Chlorophyll"
    table.cell(1, 1).text = "Green pigment"

    prs.core_properties.title = "Photosynthesis Deck"
    prs.core_properties.author = "Test Author"

    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return str(out)


def test_metadata_and_slide_count(sample_pptx):
    doc = parse_pptx(sample_pptx)
    assert doc.source.format == "pptx"
    assert doc.source.page_count == 2
    assert doc.source.title == "Photosynthesis Deck"
    assert doc.source.author == "Test Author"
    assert len(doc.pages) == 2
    assert all(page.kind == "slide" for page in doc.pages)


def test_title_heading_and_list(sample_pptx):
    doc = parse_pptx(sample_pptx)
    page1 = doc.pages[0]
    headings = [b for b in page1.blocks if b.kind == BlockKind.heading]
    assert any(h.text == "Photosynthesis" and h.level == 1 for h in headings)
    lists = [b for b in page1.blocks if b.kind == BlockKind.list]
    assert lists and lists[0].items == ["Sunlight", "Water", "Carbon dioxide"]


def test_speaker_notes_and_image(sample_pptx):
    doc = parse_pptx(sample_pptx)
    page1 = doc.pages[0]
    assert page1.notes and "three inputs" in page1.notes
    images = [b for b in page1.blocks if b.kind == BlockKind.image]
    assert images and images[0].image is not None
    assert images[0].image.width and images[0].image.height


def test_table_extracted(sample_pptx):
    doc = parse_pptx(sample_pptx)
    tables = [b for p in doc.pages for b in p.blocks if b.kind == BlockKind.table]
    assert tables
    assert tables[0].rows[0] == ["Term", "Meaning"]
    assert tables[0].rows[1] == ["Chlorophyll", "Green pigment"]


def test_json_round_trip(sample_pptx):
    doc = parse_pptx(sample_pptx)
    payload = doc.model_dump_json()
    assert '"format":"pptx"' in payload
    assert '"kind":"slide"' in payload
